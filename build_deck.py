#!/usr/bin/env python3
"""
Builds Kiro_IDE_Orchestrator_2026_Jugend.pptx — a bilingual (DE headings / EN terms)
youth-audience deck about AWS Kiro, current as of June 2026.

Run:  python3 build_deck.py
Deps: python-pptx  (pip install python-pptx)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette (dark, high-contrast, distinct cool hue per concept) ----------
BG     = RGBColor(0x0B, 0x10, 0x20)   # deep navy background
CARD   = RGBColor(0x16, 0x1F, 0x3D)   # card / panel
CARD2  = RGBColor(0x1E, 0x29, 0x4D)   # lighter card
INK    = RGBColor(0xEA, 0xEF, 0xF7)   # near-white text
MUTE   = RGBColor(0x9C, 0xA8, 0xC4)   # muted text
LINE   = RGBColor(0x2B, 0x37, 0x5E)

PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # Kiro primary
TEAL   = RGBColor(0x2D, 0xD4, 0xBF)   # requirements
BLUE   = RGBColor(0x60, 0xA5, 0xFA)   # design
INDIGO = RGBColor(0x81, 0x8C, 0xF8)   # tasks
AMBER  = RGBColor(0xFB, 0xBF, 0x24)   # "new 2026" / highlight
GREEN  = RGBColor(0x34, 0xD3, 0x99)   # good / try-it
PINK   = RGBColor(0xF4, 0x72, 0xB6)
RED    = RGBColor(0xF8, 0x71, 0x71)

FONT   = "Segoe UI"
MONO   = "Consolas"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width  = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------------------------- low-level helpers ----------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    # send to back
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s


def _set_font(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def box(s, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False, shadow=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (txt, size, color, bold, italic, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for seg in para:
            t, size, color = seg[0], seg[1], seg[2]
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            font = seg[5] if len(seg) > 5 else FONT
            r = p.add_run(); r.text = t
            _set_font(r, size, color, bold, italic, font)
    return tb


def chip(s, x, y, w, label, color, size=12.5):
    c = box(s, x, y, w, 0.42, fill=color, round_=True)
    text(s, x, y + 0.04, w, 0.34, [[(label, size, BG, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def header(s, kicker, title, accent=PURPLE):
    box(s, 0.0, 0.0, 0.18, 7.5, fill=accent)          # left accent bar
    text(s, 0.7, 0.42, 11.8, 0.4, [[(kicker, 13, accent, True)]])
    text(s, 0.7, 0.74, 11.9, 0.95, [[(title, 30, INK, True)]])
    box(s, 0.72, 1.62, 1.5, 0.045, fill=accent)
    footer(s)


def footer(s):
    text(s, 0.7, 7.06, 9.5, 0.34,
         [[("AWS Kiro · Spec-Driven Development · Stand Juni 2026", 9.5, MUTE)]])
    text(s, 10.2, 7.06, 2.45, 0.34,
         [[("kiro.dev", 9.5, MUTE)]], align=PP_ALIGN.RIGHT)


# ==================================================================== SLIDE 1
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG)
# big soft accent blocks
box(s, 8.7, -1.2, 6.5, 6.5, fill=CARD)
box(s, 9.7, 4.8, 5.0, 5.0, fill=CARD2)
chip(s, 0.9, 1.15, 2.0, "AWS · 2026", PURPLE, size=12)
text(s, 0.9, 1.95, 11.4, 2.6, [
    [("Kiro 🤖", 60, INK, True)],
    [("erst denken, dann coden.", 40, PURPLE, True)],
], line_spacing=1.0, space_after=4)
text(s, 0.92, 4.35, 9.6, 1.2, [
    [("Spec-Driven Development mit einer Agentic IDE — ", 19, INK),
     ("aus einem Satz wird ein Plan, aus dem Plan wird Code.", 19, MUTE)],
])
# three concept chips
chip(s, 0.92, 5.55, 2.7, "requirements.md", TEAL, size=12)
chip(s, 3.78, 5.55, 2.4, "design.md", BLUE, size=12)
chip(s, 6.34, 5.55, 2.2, "tasks.md", INDIGO, size=12)
text(s, 0.9, 6.5, 11.6, 0.6, [
    [("Eine Einführung für Jugendliche  ·  ", 13, MUTE),
     ("github.com/ANcpLua/Optional-Home-Exercise-Tools-Corner", 13, MUTE)],
])

# ==================================================================== SLIDE 2
s = slide()
header(s, "DAS PROBLEM", "Vibe Coding — cool, aber planlos", RED)
text(s, 0.7, 1.95, 7.0, 2.6, [
    [("Viele KI-Tools machen heute ", 16, INK), ("Vibe Coding:", 16, RED, True)],
    [("„Mach mir was Cooles\" → die KI tippt 200 Zeilen → sieht gut aus …", 15, INK)],
    [("", 6, INK)],
    [("… aber niemand weiß genau, ", 16, INK), ("was", 16, AMBER, True),
     (" das Programm können soll.", 16, INK)],
    [("Bei der ersten Änderung bricht alles zusammen — es gibt ", 15, INK),
     ("keinen Plan.", 15, RED, True)],
], line_spacing=1.12, space_after=8)
# merksatz card
box(s, 8.05, 2.0, 4.55, 3.4, fill=CARD, round_=True)
box(s, 8.05, 2.0, 0.12, 3.4, fill=RED)
text(s, 8.4, 2.3, 3.95, 3.0, [
    [("🔑 Merksatz", 15, RED, True)],
    [("", 7, INK)],
    [("Vibe Coding = plausibel", 17, INK, True)],
    [("aussehender Code, der zu", 17, INK, True)],
    [("keiner Anforderung passt.", 17, INK, True)],
    [("", 7, INK)],
    [("Klingt schlau — ist aber Raten.", 14, MUTE, False, True)],
], line_spacing=1.1, space_after=4)

# ==================================================================== SLIDE 3
s = slide()
header(s, "DIE IDEE", "Spec-Driven Development: erst der Plan", PURPLE)
text(s, 0.7, 1.78, 11.9, 0.7, [
    [("Kiro dreht die Reihenfolge um. Bevor ", 15, INK),
     ("eine einzige Zeile", 15, AMBER, True),
     (" Code entsteht, schreibt Kiro drei Dateien — ", 15, INK),
     ("„the Spec\".", 15, PURPLE, True)],
])
cards = [
    ("1  requirements.md", TEAL, "Was soll es können?", "User Stories + EARS-Notation. Die Anforderungen — testbar & eindeutig."),
    ("2  design.md", BLUE, "Wie wird es gebaut?", "Architektur, Komponenten- und Sequenz-Diagramme. Der Bauplan."),
    ("3  tasks.md", INDIGO, "Welche Schritte?", "Konkrete To-do-Liste mit Reihenfolge & Abhängigkeiten."),
]
x = 0.7
for title, col, sub, body in cards:
    box(s, x, 2.6, 3.9, 3.05, fill=CARD, round_=True)
    box(s, x, 2.6, 3.9, 0.12, fill=col)
    text(s, x + 0.28, 2.92, 3.4, 0.6, [[(title, 18, col, True, False, MONO)]])
    text(s, x + 0.28, 3.62, 3.4, 0.5, [[(sub, 15, INK, True)]])
    text(s, x + 0.28, 4.18, 3.4, 1.3, [[(body, 13.5, MUTE)]], line_spacing=1.12)
    x += 4.13
text(s, 0.7, 5.95, 11.9, 0.7, [
    [("✅ Du (der Mensch) gibst den Plan frei → dann programmiert Kiro. ", 14, GREEN, True),
     ("Jede Code-Zeile lässt sich zur Anforderung zurückverfolgen (Traceability).", 14, INK)],
])

# ==================================================================== SLIDE 4
s = slide()
header(s, "DAS HERZSTÜCK", "EARS — Anforderungen, die man testen kann", TEAL)
text(s, 0.7, 1.85, 11.9, 0.5, [
    [("Anforderungen schreibt Kiro in einem festen Muster — ", 15, INK),
     ("EARS", 15, TEAL, True), (":", 15, INK)],
])
box(s, 0.7, 2.5, 6.7, 1.5, fill=CARD, round_=True)
text(s, 1.0, 2.72, 6.1, 1.1, [
    [("WHEN ", 19, TEAL, True, False, MONO), ("[wenn das passiert]", 19, INK, False, False, MONO)],
    [("THE SYSTEM SHALL ", 19, AMBER, True, False, MONO), ("[dann tut das System das]", 19, INK, False, False, MONO)],
], line_spacing=1.25)
box(s, 7.7, 2.5, 4.9, 1.5, fill=CARD2, round_=True)
text(s, 7.95, 2.7, 4.45, 1.15, [
    [("Beispiel (Schach):", 13, TEAL, True)],
    [("WENN jemand einen verbotenen Zug macht,", 13, INK, False, False, MONO)],
    [("SOLL DAS SYSTEM den Grund anzeigen.", 13, INK, False, False, MONO)],
], line_spacing=1.12, space_after=3)
# why it matters row
items = [("🎯 Eindeutig", "kein „hab ich anders gemeint\""),
         ("🧪 Testbar", "jede Regel wird zu einem Test"),
         ("🔗 Verfolgbar", "Code ↔ Anforderung ↔ Test")]
x = 0.7
for h, b in items:
    box(s, x, 4.35, 3.9, 1.55, fill=CARD, round_=True)
    text(s, x + 0.28, 4.6, 3.4, 0.5, [[(h, 17, INK, True)]])
    text(s, x + 0.28, 5.18, 3.4, 0.6, [[(b, 13.5, MUTE)]], line_spacing=1.1)
    x += 4.13

# ==================================================================== SLIDE 5
s = slide()
header(s, "DER ABLAUF", "Vom Satz zum fertigen System", BLUE)
steps = [
    ("PROMPT", "„Bau ein Schach-\nSpiel mit FIDE-Regeln\"", PURPLE),
    ("SPEC", "Kiro schreibt die\n3 Dateien", TEAL),
    ("REVIEW", "Mensch prüft &\ngibt frei (#spec)", BLUE),
    ("IMPLEMENT", "Code — auch\nparallel (Subagents)", INDIGO),
    ("VALIDATE", "Tests grün,\nProperty-Tests", GREEN),
]
x = 0.62; w = 2.28
for i, (h, b, col) in enumerate(steps):
    box(s, x, 2.5, w, 2.4, fill=CARD, round_=True)
    box(s, x, 2.5, w, 0.12, fill=col)
    text(s, x, 2.78, w, 0.5, [[(str(i + 1), 26, col, True)]], align=PP_ALIGN.CENTER)
    text(s, x, 3.45, w, 0.5, [[(h, 14.5, INK, True)]], align=PP_ALIGN.CENTER)
    for ln_i, ln in enumerate(b.split("\n")):
        text(s, x + 0.1, 4.02 + ln_i * 0.32, w - 0.2, 0.4, [[(ln, 12, MUTE)]], align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        text(s, x + w - 0.06, 3.45, 0.4, 0.5, [[("→", 22, col, True)]], align=PP_ALIGN.CENTER)
    x += w + 0.2
text(s, 0.7, 5.45, 11.9, 0.6, [
    [("⚡ Ergebnis: ", 15, AMBER, True),
     ("Aus einer Idee wird ein nachvollziehbares, getestetes System — statt Code-Chaos.", 15, INK)],
])

# ==================================================================== SLIDE 6
s = slide()
header(s, "DIE SUPERKRÄFTE", "Was Kiro besonders macht", INDIGO)
feats = [
    ("🪝 Agent Hooks", TEAL, "„Immer wenn ich speichere, aktualisiere die Docs.\" Du beschreibst es, Kiro baut den Hook."),
    ("🧭 Steering Files", BLUE, "Kiros Hausregeln: Stil & Namens-Konventionen in .kiro/steering/ — alles bleibt einheitlich."),
    ("🖼️ Multimodal", PINK, "Foto einer Whiteboard-Skizze hochladen → Kiro macht Architektur & Code daraus."),
    ("🧩 MCP", AMBER, "Verbindet externe Tools: GitHub, Docs, eigene CLI-Tools — über das Model Context Protocol."),
    ("👥 Subagents", INDIGO, "Mehrere Aufgaben gleichzeitig: spezialisierte Agenten arbeiten parallel, jeder mit eigenem Kontext."),
    ("🔐 Permissions", GREEN, "IDE 1.0 (Juni 2026): du bestimmst genau, was die KI tun darf — Agent-Focus & Custom Agents."),
]
positions = [(0.7, 1.95), (4.83, 1.95), (8.96, 1.95),
             (0.7, 4.4), (4.83, 4.4), (8.96, 4.4)]
for (px, py), (h, col, b) in zip(positions, feats):
    box(s, px, py, 3.67, 2.25, fill=CARD, round_=True)
    box(s, px, py, 0.1, 2.25, fill=col)
    text(s, px + 0.28, py + 0.22, 3.2, 0.5, [[(h, 16, INK, True)]])
    text(s, px + 0.28, py + 0.82, 3.25, 1.3, [[(b, 12.5, MUTE)]], line_spacing=1.1)

# ==================================================================== SLIDE 7
s = slide()
header(s, "WICHTIG FÜRS PRÄSENTIEREN", "Was 2026 neu ist", AMBER)
tl = [
    ("Jul 2025", "Public Preview", PURPLE),
    ("17. Nov 2025", "Generally Available (GA) — kein Beta mehr", GREEN),
    ("Feb 2026", "Subagents, Skills & Hooks · Okta / Microsoft Entra ID SSO", BLUE),
    ("Mai 2026", "Concurrent Spec-Tasks · Quick-Plan-Workflow", TEAL),
    ("25. Jun 2026", "IDE 1.0 — Agent Focus, Permissions, Custom Agents", AMBER),
]
y = 1.95
for date, label, col in tl:
    box(s, 0.95, y + 0.06, 0.22, 0.22, fill=col, round_=True)
    text(s, 1.35, y, 2.3, 0.4, [[(date, 14.5, col, True)]])
    text(s, 3.75, y, 8.7, 0.4, [[(label, 14.5, INK)]])
    if (date, label, col) != tl[-1]:
        box(s, 1.04, y + 0.34, 0.04, 0.46, fill=LINE)
    y += 0.82
box(s, 0.7, 6.18, 11.9, 0.7, fill=CARD, round_=True)
text(s, 1.0, 6.32, 11.3, 0.5, [
    [("Neu außerdem: ", 13.5, AMBER, True),
     ("Kiro CLI (Terminal) · Kiro Web (app.kiro.dev) · Web-Suche im Chat · „Auto\"-Modell-Agent", 13.5, INK)],
])

# ==================================================================== SLIDE 8
s = slide()
header(s, "PREISE", "Was kostet das? (Stand Juni 2026)", GREEN)
text(s, 0.7, 1.8, 11.9, 0.5, [
    [("Kiro rechnet mit ", 15, INK), ("Credits", 15, GREEN, True),
     (" — 1 Credit = ein Stück Arbeit; einfache Prompts kosten weniger als 1 Credit.", 15, INK)],
])
rows = [("Free", "0 $", "50", GREEN),
        ("Pro", "20 $", "1.000", BLUE),
        ("Pro+", "40 $", "2.000", INDIGO),
        ("Pro Max", "100 $", "5.000", PURPLE),
        ("Power", "200 $", "10.000", AMBER)]
# table header
hx = [0.7, 4.9, 8.0]; hw = [4.0, 2.9, 4.0]
heads = ["Plan", "Preis / Monat", "Credits / Monat"]
box(s, 0.7, 2.45, 11.9, 0.5, fill=CARD2, round_=True)
for hxi, hwi, ht in zip(hx, hw, heads):
    text(s, hxi + 0.2, 2.5, hwi, 0.4, [[(ht, 13.5, MUTE, True)]], anchor=MSO_ANCHOR.MIDDLE)
y = 3.02
for plan, price, cr, col in rows:
    box(s, 0.7, y, 11.9, 0.56, fill=CARD, round_=True)
    box(s, 0.7, y, 0.1, 0.56, fill=col)
    text(s, 0.95, y, 3.8, 0.56, [[(plan, 15.5, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.1, y, 2.7, 0.56, [[(price, 15, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.2, y, 3.8, 0.56, [[(cr, 15, INK)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.63
text(s, 0.7, 6.55, 11.9, 0.5, [
    [("👉 Zum Ausprobieren reicht ", 13, GREEN, True), ("Free", 13, GREEN, True),
     (" (Social-Login / AWS Builder ID — kein Kreditkarten-Konto nötig). Overage: $0,04 / Credit. Preise: kiro.dev/pricing", 13, MUTE)],
])

# ==================================================================== SLIDE 9
s = slide()
header(s, "SELBST AUSPROBIEREN", "In 5 Schritten loslegen", TEAL)
steps = [
    ("1", "Download", "kiro.dev — ein VS-Code-Fork, sieht vertraut aus.", PURPLE),
    ("2", "Anmelden", "Free-Plan via Social-Login. Kein Amazon-Konto nötig.", BLUE),
    ("3", "Spec-Modus", "Wunsch eintippen: „Build a to-do app with a checkbox per item.\"", TEAL),
    ("4", "Plan prüfen", "requirements.md · design.md · tasks.md anschauen & freigeben.", INDIGO),
    ("5", "Implementieren", "Kiro coden lassen — und zusehen, wie alles zur Spec passt.", GREEN),
]
y = 1.95
for num, h, b, col in steps:
    box(s, 0.7, y, 11.9, 0.92, fill=CARD, round_=True)
    box(s, 0.7, y, 0.92, 0.92, fill=col, round_=True)
    text(s, 0.7, y, 0.92, 0.92, [[(num, 30, BG, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.85, y + 0.12, 3.0, 0.7, [[(h, 17, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.7, y, 7.7, 0.92, [[(b, 14, MUTE)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 1.02

# ==================================================================== SLIDE 10
s = slide()
header(s, "VERGLEICH", "Kiro vs. „normale\" KI-Tools", PINK)
cols_x = [0.7, 5.1, 8.85]
box(s, 5.1, 1.95, 3.55, 4.6, fill=CARD2, round_=True)      # highlight Kiro column
box(s, 8.85, 1.95, 3.75, 4.6, fill=CARD, round_=True)
# headers
text(s, 0.7, 2.1, 4.2, 0.6, [[("", 14, INK)]])
text(s, 5.1, 2.12, 3.55, 0.6, [[("Vibe Coding", 18, RED, True)]], align=PP_ALIGN.CENTER)
text(s, 8.85, 2.12, 3.75, 0.6, [[("Kiro (Spec-Driven)", 18, GREEN, True)]], align=PP_ALIGN.CENTER)
rows = [
    ("Reihenfolge", "Sofort Code", "Erst Plan, dann Code"),
    ("Nachvollziehbar?", "Oft nicht", "Ja — Zeile → Anforderung"),
    ("Bei Änderungen", "Bricht schnell", "Plan passt sich mit an"),
    ("Gut für", "Schnelle Experimente", "Echte, wartbare Projekte"),
]
y = 2.95
for label, vibe, kiro in rows:
    text(s, 0.9, y, 4.0, 0.5, [[(label, 14.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.1, y, 3.55, 0.5, [[(vibe, 13.5, MUTE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 8.85, y, 3.75, 0.5, [[(kiro, 13.5, INK)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if (label, vibe, kiro) != rows[-1]:
        box(s, 0.9, y + 0.62, 11.5, 0.012, fill=LINE)
    y += 0.86

# ==================================================================== SLIDE 11
s = slide()
box(s, 0, 0, 13.333, 7.5, fill=BG)
box(s, 0, 0, 13.333, 0.18, fill=PURPLE)
box(s, 0, 7.32, 13.333, 0.18, fill=PURPLE)
text(s, 0.9, 1.5, 11.6, 1.2, [[("Das Take-away 🏁", 40, INK, True)]])
box(s, 0.9, 2.85, 11.55, 1.9, fill=CARD, round_=True)
box(s, 0.9, 2.85, 0.14, 1.9, fill=AMBER)
text(s, 1.35, 3.15, 11.0, 1.4, [
    [("Kiro = „erst denken, dann coden\".", 26, AMBER, True)],
    [("Aus einem Satz wird ein ", 17, INK), ("Plan (Spec)", 17, TEAL, True),
     (", aus dem Plan wird ", 17, INK), ("nachvollziehbarer Code.", 17, GREEN, True)],
    [("Das ist der Unterschied zwischen ", 16, MUTE), ("Raten", 16, RED, True),
     (" und ", 16, MUTE), ("Engineering.", 16, INK, True)],
], line_spacing=1.15, space_after=8)
text(s, 0.9, 5.25, 11.6, 0.9, [
    [("Mehr:  ", 15, MUTE),
     ("kiro.dev", 15, BLUE, True), ("   ·   ", 15, MUTE),
     ("kiro.dev/docs", 15, BLUE, True), ("   ·   ", 15, MUTE),
     ("kiro.dev/changelog/ide", 15, BLUE, True)],
])
text(s, 0.9, 6.15, 11.6, 0.6, [
    [("Danke! Fragen? 🙋  ", 16, INK, True),
     ("— probiert den Free-Plan aus und baut euer erstes Spec.", 16, MUTE)],
])

prs.save("Kiro_IDE_Orchestrator_2026_Jugend.pptx")
print("OK: Kiro_IDE_Orchestrator_2026_Jugend.pptx  (%d slides)" % len(prs.slides._sldIdLst))
